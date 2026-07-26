"""Slide deck ingestion — one chunk per slide from PDF or PPTX files.

Decks are visual-first content: each slide is its own semantic unit, so we
emit one chunk per slide (no sub-splitting) with the 1-indexed slide number.
Like papers, there is no time dimension — t_start/t_end are always 0.0.

PDF decks: each page is treated as one slide (pymupdf text extraction).
PPTX decks: text from every shape on the slide is concatenated.
Slides with negligible text (< 10 chars) are skipped — they're typically
title cards or full-bleed images that don't contribute searchable content.
"""
from __future__ import annotations

from pathlib import Path

_MIN_SLIDE_CHARS = 10


def _parse_pdf_deck(path: Path) -> list[dict]:
    """Extract one chunk per page from a PDF slide deck."""
    import fitz  # pymupdf

    try:
        doc = fitz.open(str(path))
    except Exception as exc:
        print(f"[deck] failed to open {path.name} ({type(exc).__name__}: {exc})")
        return []

    chunks: list[dict] = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        try:
            text = (page.get_text() or "").strip()
        except Exception:
            text = ""

        if len(text) < _MIN_SLIDE_CHARS:
            continue

        chunks.append({
            "text": text,
            "slide": page_num + 1,
            "t_start": 0.0,
            "t_end": 0.0,
        })

    doc.close()
    return chunks


def _parse_pptx_deck(path: Path) -> list[dict]:
    """Extract one chunk per slide from a PPTX file."""
    from pptx import Presentation

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        print(f"[deck] failed to open {path.name} ({type(exc).__name__}: {exc})")
        return []

    chunks: list[dict] = []
    for slide_idx, slide in enumerate(prs.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        parts.append(para_text)

        text = "\n".join(parts).strip()
        if len(text) < _MIN_SLIDE_CHARS:
            continue

        chunks.append({
            "text": text,
            "slide": slide_idx + 1,
            "t_start": 0.0,
            "t_end": 0.0,
        })

    return chunks


def parse_deck(path: Path) -> list[dict]:
    """Parse a slide deck (PDF or PPTX) into one chunk per slide.

    Dispatches to the right backend based on file extension. Slides with
    fewer than 10 characters of text are skipped.
    """
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return _parse_pptx_deck(path)
    if suffix == ".pdf":
        return _parse_pdf_deck(path)
    print(f"[deck] unsupported format: {suffix}")
    return []
